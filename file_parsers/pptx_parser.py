"""
PPTX parser — extracts text from PowerPoint presentations using python-pptx.
"""

import io
import logging
from .base import BaseParser

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Extract text from .pptx files using python-pptx."""

    supported_mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]

    def parse(self, data: bytes, filename: str = "") -> dict:
        try:
            from pptx import Presentation  # python-pptx
        except ImportError:
            raise RuntimeError(
                "python-pptx is not installed. Run: pip install python-pptx"
            )

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            logger.error(f"PPTX parsing failed for '{filename}': {exc}")
            raise ValueError(f"Failed to parse PPTX: {exc}") from exc

        slides_text: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        slide_parts.append(line)
            if slide_parts:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)
        return {
            "text": full_text,
            "slide_count": len(prs.slides),
        }
